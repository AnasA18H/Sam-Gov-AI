"""
Text Extraction Service
Extracts raw text from ALL common document formats using free/open-source libraries:
- PDF (text-based and scanned with OCR)
- Microsoft Word (.docx, .doc)
- Microsoft Excel (.xlsx, .xls, .csv)
- Images (.jpg, .png, .tiff, .bmp) with OCR
- PowerPoint (.pptx)
- Plain Text (.txt, .rtf, .md)
"""
import re
import logging
import mimetypes
import struct
from pathlib import Path
from typing import Any, Optional, Tuple
import io

# Document processing libraries
import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook

# Additional format support (optional imports)
try:
    import pandas as pd
    _pd = pd
    PANDAS_AVAILABLE = True
except ImportError:
    _pd = None  # type: ignore[assignment]
    PANDAS_AVAILABLE = False
    logging.warning("pandas not available. CSV and advanced Excel features disabled.")

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False
    logging.debug("xlrd not available. .xls files will use fallback methods.")

try:
    from pptx import Presentation as _Presentation
    PPTX_AVAILABLE = True
except ImportError:
    _Presentation = None  # type: ignore[assignment]
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint support disabled.")

try:
    import pytesseract as _pytesseract
    from PIL import Image as _PilImage
    import cv2 as _cv2
    import numpy as _np
    # Configure pytesseract to use system tesseract
    _pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
    OCR_AVAILABLE = True
except ImportError as e:
    _pytesseract = None  # type: ignore[assignment]
    _PilImage = None  # type: ignore[assignment]
    _cv2 = None  # type: ignore[assignment]
    _np = None  # type: ignore[assignment]
    OCR_AVAILABLE = False
    logging.warning(f"OCR libraries (pytesseract, PIL, cv2) not available: {str(e)}. Scanned PDF and image OCR disabled.")
except Exception as e:
    _pytesseract = None  # type: ignore[assignment]
    _PilImage = None  # type: ignore[assignment]
    _cv2 = None  # type: ignore[assignment]
    _np = None  # type: ignore[assignment]
    OCR_AVAILABLE = False
    logging.warning(f"OCR configuration error: {str(e)}. OCR disabled.")

try:
    import PyPDF2 as _PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    _PyPDF2 = None  # type: ignore[assignment]
    PYPDF2_AVAILABLE = False
    logging.debug("PyPDF2 not available. Using pdfplumber only for PDFs.")

try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    _rtf_to_text = None  # type: ignore[assignment]
    RTF_AVAILABLE = False
    logging.debug("striprtf not available. RTF files will use basic text extraction.")

try:
    import markdown as _markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    _markdown = None  # type: ignore[assignment]
    MARKDOWN_AVAILABLE = False
    logging.debug("markdown not available. Markdown files will use basic text extraction.")

# Google Document AI (optional - for high-quality text extraction)
try:
    from google.cloud import documentai as _documentai
    from google.oauth2 import service_account as _service_account
    DOCAI_AVAILABLE = True
except ImportError:
    _documentai = None  # type: ignore[assignment]
    _service_account = None  # type: ignore[assignment]
    DOCAI_AVAILABLE = False
    logging.debug("google-cloud-documentai not available. Document AI extraction disabled.")
except Exception as e:
    _documentai = None  # type: ignore[assignment]
    _service_account = None  # type: ignore[assignment]
    DOCAI_AVAILABLE = False
    logging.debug(f"Document AI initialization error: {str(e)}. Document AI extraction disabled.")

from ..core.config import settings

logger = logging.getLogger(__name__)


class TextExtractor:
    """Service for extracting raw text from various document formats"""
    
    # Document type classification patterns
    DOCUMENT_TYPE_PATTERNS = {
        'sf1449': [r'SF\s*1449', r'Standard Form\s*1449', r'Form\s*1449', r'Section B - Schedule of Supplies/Services'],
        'sf33': [r'SF\s*33', r'Standard Form\s*33', r'Solicitation, Offer and Award', r'SECTION B - SUPPLIES OR SERVICES AND PRICES/COSTS'],
        'sf18': [r'SF\s*18', r'Request for Quotation', r'Block 11 - Schedule'],
        'sf30': [r'SF\s*30', r'Standard Form\s*30', r'Amendment', r'Modification'],
        'dla': [r'DLA', r'DIBBS', r'PR:', r'PRLI', r'NSN/MATERIAL:'],
        'va': [r'Department of Veterans Affairs', r'VA Solicitation', r'Price/Cost Schedule'],
        'gsa': [r'GSA Schedule', r'GSA Advantage', r'Price List'],
        'sow': [r'Statement\s+of\s+Work', r'SOW', r'Scope\s+of\s+Work'],
        'amendment': [r'Amendment', r'Modification', r'Change\s+Order'],
    }
    
    def __init__(self):
        """Initialize the text extractor"""
        # Initialize OCR (pytesseract configuration)
        if OCR_AVAILABLE and _pytesseract is not None:
            try:
                # Ensure tesseract path is set
                if not hasattr(_pytesseract.pytesseract, 'tesseract_cmd') or not _pytesseract.pytesseract.tesseract_cmd:
                    _pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
                logger.info("OCR (pytesseract) initialized successfully")
            except Exception as e:
                logger.warning(f"OCR initialization warning: {str(e)}")
    
    @staticmethod
    def _clean_text(text: str, preserve_layout: bool = True) -> str:
        """
        Clean extracted text by removing PDF encoding artifacts and other garbage
        Removes (cid:XXX) patterns which are PDF character encoding artifacts
        Also removes corrupted patterns like semicolon-separated garbage text
        """
        if not text:
            return ""
        
        # Remove PDF encoding artifacts (cid:XXX) where XXX is any number
        text = re.sub(r'\(cid:\d+\)', '', text)
        
        # Remove corrupted patterns: semicolon-separated alphanumeric garbage
        # Pattern: ; followed by alphanumeric chars, <, >, =, :, etc. ending with ;
        # Example: ;9<746Q57BG69B<;4C8=B51A75BA74:5;38<=5914:A==93A;
        text = re.sub(r';[A-Z0-9<>=:]+(?:;[A-Z0-9<>=:]+)+;', '', text)
        
        # Remove lines that are mostly garbage (high ratio of special chars to letters)
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # Skip lines that are mostly special characters and numbers with few actual words
            if len(line.strip()) > 0:
                # Count actual letters (a-z, A-Z)
                letter_count = len(re.findall(r'[a-zA-Z]', line))
                # Count total alphanumeric chars
                alnum_count = len(re.findall(r'[a-zA-Z0-9]', line))
                total_chars = len(line.strip())
                
                # If line has very few letters relative to total chars, or is mostly special chars
                # BUT: Be less aggressive - keep lines that might be CLIN tables (have numbers, CLIN patterns, etc.)
                # Only skip if it's clearly garbage (no letters AND no meaningful numbers/patterns)
                has_clin_pattern = bool(re.search(r'(CLIN|Item\s*(?:No|Number)|Line\s*Item|000\d|NSN)', line, re.IGNORECASE))
                has_meaningful_numbers = bool(re.search(r'\d{3,}', line))  # At least 3 consecutive digits
                
                if total_chars > 10 and (letter_count < 3 or (alnum_count / total_chars) < 0.3):
                    # Only skip if it doesn't have CLIN patterns or meaningful numbers
                    if not has_clin_pattern and not has_meaningful_numbers:
                        # This is likely garbage, skip it
                        continue
            cleaned_lines.append(line)
        text = '\n'.join(cleaned_lines)
        
        # Final cleaning
        if preserve_layout:
            # Preserve structure: remove non-printable but keep newlines/tabs
            text = "".join(c for c in text if ord(c) >= 32 or c in '\n\r\t')
            # Normalize excessive empty lines but keep single/double newlines
            text = re.sub(r'\n{3,}', '\n\n', text)
            # Remove trailing spaces from each line to save tokens, but keep leading/internal spacing
            text = '\n'.join(line.rstrip() for line in text.split('\n'))
        else:
            # Traditional cleaning: collapse all whitespace including newlines
            text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)
            text = re.sub(r'\s+', ' ', text)
        
        return text.strip()
    
    def classify_document_type(self, file_path: Path, text: str) -> str:
        """
        Classify document type to route to correct extraction method
        
        Returns:
            Document type: 'sf1449', 'sf30', 'sow', 'amendment', or 'unknown'
        """
        # Check filename first
        filename_lower = file_path.name.lower()
        for doc_type, patterns in self.DOCUMENT_TYPE_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, filename_lower, re.IGNORECASE):
                    logger.info(f"Document classified as {doc_type} based on filename")
                    return doc_type
        
        # Check text content
        text_lower = text[:2000].lower()  # Check first 2000 chars
        for doc_type, patterns in self.DOCUMENT_TYPE_PATTERNS.items():
            for pattern in patterns:
                if re.search(pattern, text_lower, re.IGNORECASE):
                    logger.info(f"Document classified as {doc_type} based on content")
                    return doc_type
        
        return 'unknown'
    
    @staticmethod
    def _detect_file_format(file_path: Path) -> Tuple[str, str]:
        """
        Detect file format using both extension and magic bytes.
        
        Returns:
            Tuple of (format_type, mime_type)
        """
        file_ext = file_path.suffix.lower()
        
        # Read magic bytes (first few bytes of file)
        try:
            with open(file_path, 'rb') as f:
                magic_bytes = f.read(16)
        except Exception:
            magic_bytes = b''
        
        # PDF detection
        if file_ext == '.pdf' or magic_bytes.startswith(b'%PDF'):
            return 'pdf', 'application/pdf'
        
        # Word documents
        if file_ext == '.docx' or (magic_bytes.startswith(b'PK') and b'word/' in magic_bytes[:1000]):
            return 'docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        if file_ext == '.doc' or magic_bytes.startswith(b'\xd0\xcf\x11\xe0'):
            return 'doc', 'application/msword'
        
        # Excel documents
        if file_ext == '.xlsx' or (magic_bytes.startswith(b'PK') and b'xl/' in magic_bytes[:1000]):
            return 'xlsx', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        if file_ext == '.xls' or magic_bytes.startswith(b'\xd0\xcf\x11\xe0'):
            return 'xls', 'application/vnd.ms-excel'
        if file_ext == '.csv':
            return 'csv', 'text/csv'
        
        # PowerPoint
        if file_ext == '.pptx' or (magic_bytes.startswith(b'PK') and b'ppt/' in magic_bytes[:1000]):
            return 'pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        # Images
        if file_ext in ['.jpg', '.jpeg'] or magic_bytes.startswith(b'\xff\xd8\xff'):
            return 'image', 'image/jpeg'
        if file_ext == '.png' or magic_bytes.startswith(b'\x89PNG'):
            return 'image', 'image/png'
        if file_ext == '.tiff' or magic_bytes.startswith(b'II*\x00') or magic_bytes.startswith(b'MM\x00*'):
            return 'image', 'image/tiff'
        if file_ext == '.bmp' or magic_bytes.startswith(b'BM'):
            return 'image', 'image/bmp'
        
        # Text formats
        if file_ext == '.txt':
            return 'txt', 'text/plain'
        if file_ext == '.rtf' or magic_bytes.startswith(b'{\\rtf'):
            return 'rtf', 'application/rtf'
        if file_ext == '.md':
            return 'md', 'text/markdown'
        
        # Default based on extension
        mime_type, _ = mimetypes.guess_type(str(file_path))
        return file_ext[1:] if file_ext else 'unknown', mime_type or 'application/octet-stream'
    
    def extract_text(self, file_path: str) -> str:
        """
        Extract text from a document file with automatic format detection.
        Supports PDF, Word, Excel, PowerPoint, Images, and Text formats.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content
        """
        file_path_obj = Path(file_path)
        
        # Handle relative paths (for Digital Ocean compatibility)
        if not file_path_obj.is_absolute():
            project_root = Path(settings.PROJECT_ROOT)
            abs_path = project_root / file_path
            if abs_path.exists():
                file_path_obj = abs_path
            elif hasattr(settings, 'STORAGE_BASE_PATH'):
                storage_base = Path(settings.STORAGE_BASE_PATH)
                abs_path = storage_base.parent / file_path if 'backend/data' in file_path else storage_base / file_path
                if abs_path.exists():
                    file_path_obj = abs_path
        
        if not file_path_obj.exists():
            raise FileNotFoundError(f"Document file not found: {file_path}")
        
        # Detect file format
        format_type, mime_type = self._detect_file_format(file_path_obj)
        logger.info(f"Detected format: {format_type} (MIME: {mime_type}) for {file_path_obj.name}")
        
        try:
            # Route to appropriate extraction method
            if format_type == 'pdf':
                return self._extract_from_pdf(file_path_obj)
            elif format_type in ['docx', 'doc']:
                return self._extract_from_word(file_path_obj, format_type)
            elif format_type in ['xlsx', 'xls', 'csv']:
                return self._extract_from_excel(file_path_obj, format_type)
            elif format_type == 'pptx':
                return self._extract_from_powerpoint(file_path_obj)
            elif format_type == 'image':
                return self._extract_from_image(file_path_obj)
            elif format_type == 'txt':
                return self._extract_from_text(file_path_obj)
            elif format_type == 'rtf':
                return self._extract_from_rtf(file_path_obj)
            elif format_type == 'md':
                return self._extract_from_markdown(file_path_obj)
            else:
                logger.warning(f"Unsupported file format: {format_type} for {file_path_obj.name}")
                # Try as plain text as fallback
                try:
                    return self._extract_from_text(file_path_obj)
                except Exception:
                    return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}", exc_info=True)
            return ""

    def extract_text_pdf_with_document_ai(self, file_path: str | Path) -> str:
        """
        Extract text from a PDF using Document AI OCR only (no pdfplumber).
        Use this when you want consistent OCR output for static forms (e.g. KVP extraction).
        Handles large PDFs via 15-page chunking. Requires GOOGLE_* config and credentials.

        Args:
            file_path: Path to PDF file

        Returns:
            Extracted text from Document AI

        Raises:
            FileNotFoundError: If file missing or not PDF
            ImportError/Exception: If Document AI unavailable or processing fails
        """
        file_path_obj = Path(file_path)
        if not file_path_obj.is_absolute():
            project_root = Path(settings.PROJECT_ROOT)
            abs_path = project_root / file_path_obj
            if abs_path.exists():
                file_path_obj = abs_path
            elif hasattr(settings, 'STORAGE_BASE_PATH'):
                storage_base = Path(settings.STORAGE_BASE_PATH)
                abs_path = storage_base.parent / str(file_path_obj) if 'backend/data' in str(file_path_obj) else storage_base / file_path_obj
                if abs_path.exists():
                    file_path_obj = abs_path
        if not file_path_obj.exists():
            raise FileNotFoundError(f"PDF file not found: {file_path_obj}")
        format_type, _ = self._detect_file_format(file_path_obj)
        if format_type != 'pdf':
            raise ValueError(f"Not a PDF: {file_path_obj.name} (detected: {format_type})")
        return self._extract_with_document_ai(file_path_obj)

    def _extract_with_document_ai(self, file_path: Path) -> str:
        """
        Extract text from PDF using Google Document AI.
        Handles large PDFs by splitting into chunks if needed.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text content
        """
        if not DOCAI_AVAILABLE:
            raise ImportError("google-cloud-documentai not available")
        
        # Try to find service account JSON file
        json_path = Path(settings.GOOGLE_SERVICE_ACCOUNT_JSON) if settings.GOOGLE_SERVICE_ACCOUNT_JSON else None
        
        # If not set, try default location
        if not json_path or not json_path.exists():
            default_path = settings.PROJECT_ROOT / "extras" / "resolute-planet-485419-f8-f543cf0a64b5.json"
            if default_path.exists():
                json_path = default_path
            else:
                raise FileNotFoundError(f"Service account JSON not found. Set GOOGLE_SERVICE_ACCOUNT_JSON or place file at: {default_path}")
        
        try:
            if _service_account is None or _documentai is None:
                raise ImportError("google-cloud-documentai not available")
            # Authenticate
            creds = _service_account.Credentials.from_service_account_file(
                str(json_path),
                scopes=["https://www.googleapis.com/auth/cloud-platform"]
            )
            
            # Create client
            client = _documentai.DocumentProcessorServiceClient(credentials=creds)
            
            # Check PDF page count (Document AI has 15 page limit per request)
            page_count = None
            try:
                if PYPDF2_AVAILABLE and _PyPDF2 is not None:
                    with open(file_path, "rb") as f:
                        pdf_reader = _PyPDF2.PdfReader(f)
                        page_count = len(pdf_reader.pages)
            except Exception:
                pass
            
            # Process document (split if needed)
            if page_count and page_count > 15:
                logger.info(f"PDF has {page_count} pages, splitting into chunks for Document AI")
                return self._extract_large_pdf_with_docai(file_path, client, page_count)
            else:
                # Process normally
                with open(file_path, "rb") as f:
                    file_content = f.read()
                
                name = f"projects/{settings.GOOGLE_PROJECT_ID}/locations/{settings.GOOGLE_LOCATION}/processors/{settings.GOOGLE_PROCESSOR_ID}"
                
                raw_document = _documentai.RawDocument(
                    content=file_content,
                    mime_type="application/pdf"
                )
                
                request = _documentai.ProcessRequest(
                    name=name,
                    raw_document=raw_document
                )
                
                logger.debug(f"Processing PDF with Document AI: {file_path.name}")
                response = client.process_document(request=request)
                
                text = response.document.text
                logger.info(f"Document AI extracted {len(text)} characters from {len(response.document.pages)} pages")
                return text

        except Exception as e:
            logger.error(f"Document AI extraction failed: {str(e)}", exc_info=True)
            raise

    def get_document_ai_document(self, file_path: str | Path):
        """
        Call Document AI and return the full Document object (with pages/lines/layout)
        for OCR position extraction. Returns None if PDF has >15 pages (chunked) or on error.
        """
        file_path_obj = Path(file_path)
        if not file_path_obj.is_absolute():
            project_root = Path(settings.PROJECT_ROOT)
            abs_path = project_root / file_path_obj
            if abs_path.exists():
                file_path_obj = abs_path
        if not file_path_obj.exists():
            raise FileNotFoundError(f"PDF file not found: {file_path_obj}")
        format_type, _ = self._detect_file_format(file_path_obj)
        if format_type != 'pdf':
            raise ValueError(f"Not a PDF: {file_path_obj.name}")
        if not DOCAI_AVAILABLE:
            return None
        page_count = None
        try:
            if PYPDF2_AVAILABLE and _PyPDF2 is not None:
                with open(file_path_obj, "rb") as f:
                    pdf_reader = _PyPDF2.PdfReader(f)
                    page_count = len(pdf_reader.pages)
        except Exception:
            pass
        if page_count and page_count > 15:
            logger.info("ocr_fill: get_document_ai_document skipped (PDF has %s pages, positions only for <=15)", page_count)
            return None
        try:
            json_path = Path(settings.GOOGLE_SERVICE_ACCOUNT_JSON) if settings.GOOGLE_SERVICE_ACCOUNT_JSON else None
            if not json_path or not json_path.exists():
                default_path = settings.PROJECT_ROOT / "extras" / "resolute-planet-485419-f8-f543cf0a64b5.json"
                json_path = default_path if default_path.exists() else None
            if not json_path or not json_path.exists():
                return None
            if _service_account is None or _documentai is None:
                return None
            creds = _service_account.Credentials.from_service_account_file(str(json_path), scopes=["https://www.googleapis.com/auth/cloud-platform"])
            client = _documentai.DocumentProcessorServiceClient(credentials=creds)
            name = f"projects/{settings.GOOGLE_PROJECT_ID}/locations/{settings.GOOGLE_LOCATION}/processors/{settings.GOOGLE_PROCESSOR_ID}"
            with open(file_path_obj, "rb") as f:
                file_content = f.read()
            raw_document = _documentai.RawDocument(content=file_content, mime_type="application/pdf")
            response = client.process_document(request=_documentai.ProcessRequest(name=name, raw_document=raw_document))
            return response.document
        except Exception as e:
            logger.warning("ocr_fill: get_document_ai_document failed: %s", e)
            return None
    
    def _extract_large_pdf_with_docai(self, file_path: Path, client, total_pages: int) -> str:
        """
        Process large PDFs by splitting into 15-page chunks (Document AI limit).
        
        Args:
            file_path: Path to PDF file
            client: Document AI client
            total_pages: Total number of pages in PDF
            
        Returns:
            Combined extracted text from all chunks
        """
        if not PYPDF2_AVAILABLE or _PyPDF2 is None or _documentai is None:
            raise ImportError("PyPDF2 and Document AI required for splitting large PDFs")

        name = f"projects/{settings.GOOGLE_PROJECT_ID}/locations/{settings.GOOGLE_LOCATION}/processors/{settings.GOOGLE_PROCESSOR_ID}"
        all_text_parts = []
        chunk_size = 15  # Document AI limit

        try:
            with open(file_path, "rb") as f:
                pdf_content = f.read()

            pdf_reader = _PyPDF2.PdfReader(io.BytesIO(pdf_content))
            total_chunks = (total_pages + chunk_size - 1) // chunk_size
            
            logger.info(f"Processing {total_chunks} chunk(s) of up to {chunk_size} pages each...")
            
            for chunk_num in range(total_chunks):
                start_page = chunk_num * chunk_size
                end_page = min(start_page + chunk_size, total_pages)
                
                # Create a new PDF with just these pages
                pdf_writer = _PyPDF2.PdfWriter()
                for page_num in range(start_page, end_page):
                    pdf_writer.add_page(pdf_reader.pages[page_num])
                
                # Write chunk to bytes
                chunk_buffer = io.BytesIO()
                pdf_writer.write(chunk_buffer)
                chunk_content = chunk_buffer.getvalue()
                
                logger.debug(f"Processing chunk {chunk_num + 1}/{total_chunks} (pages {start_page + 1}-{end_page})...")
                
                # Process chunk
                raw_document = _documentai.RawDocument(
                    content=chunk_content,
                    mime_type="application/pdf"
                )
                
                request = _documentai.ProcessRequest(
                    name=name,
                    raw_document=raw_document
                )
                
                response = client.process_document(request=request)
                chunk_text = response.document.text
                all_text_parts.append(chunk_text)
            
            # Combine all chunks
            combined_text = "\n\n--- Page Break ---\n\n".join(all_text_parts)
            logger.info(f"Document AI extracted {len(combined_text)} characters from {total_pages} pages (split into {total_chunks} chunks)")
            return combined_text
            
        except Exception as e:
            logger.error(f"Error processing large PDF with Document AI: {str(e)}", exc_info=True)
            raise
    
    def _detect_pdf_type(self, file_path: Path) -> str:
        """
        Detect if PDF is text-based or scanned/image-based.
        Uses multiple heuristics to determine PDF type accurately.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            'text_based' or 'scanned'
        """
        try:
            with pdfplumber.open(file_path) as pdf:
                # Sample first few pages to determine type
                sample_pages = min(3, len(pdf.pages))
                total_text_chars = 0
                total_pages_checked = 0
                pages_with_text_objects = 0
                
                for page_num in range(sample_pages):
                    page = pdf.pages[page_num]
                    
                    # Check 1: Try to extract text directly
                    page_text = page.extract_text()
                    if page_text:
                        text_length = len(page_text.strip())
                        total_text_chars += text_length
                        total_pages_checked += 1
                        
                        # If we get substantial text (>200 chars), likely text-based
                        if text_length > 200:
                            logger.debug(f"PDF appears text-based: page {page_num + 1} has {text_length} chars")
                            return 'text_based'
                    
                    # Check 2: Look for text objects in PDF structure (most reliable indicator)
                    try:
                        chars = page.chars
                        if chars and len(chars) > 100:  # Need substantial text objects
                            pages_with_text_objects += 1
                            logger.debug(f"Page {page_num + 1} has {len(chars)} text objects")
                    except:
                        pass
                    
                    # Check 3: Look for images (scanned PDFs often have large images)
                    try:
                        images = page.images
                        if images and len(images) > 0:
                            # Check if images are large (likely scanned pages)
                            large_images = [img for img in images if img.get('width', 0) > 500 and img.get('height', 0) > 500]
                            if large_images and not page_text:
                                logger.debug(f"Page {page_num + 1} has large images and no text - likely scanned")
                                return 'scanned'
                    except:
                        pass
                
                # Decision logic
                # If most pages have text objects, it's text-based
                if pages_with_text_objects >= sample_pages * 0.7:
                    logger.debug(f"PDF appears text-based: {pages_with_text_objects}/{sample_pages} pages have text objects")
                    return 'text_based'
                
                # If average text per page is very low, likely scanned
                if total_pages_checked > 0:
                    avg_text = total_text_chars / total_pages_checked
                    if avg_text < 100:  # Lower threshold for scanned detection
                        logger.debug(f"PDF appears scanned: average {avg_text:.1f} chars per page")
                        return 'scanned'
                
                # If we got some text but not much, check quality
                # Text-based PDFs usually have clean, structured text
                # Scanned PDFs extracted with pdfplumber often have poor quality
                if total_text_chars > 0 and total_text_chars < 500:
                    # Try to check if text looks like OCR output (has artifacts)
                    sample_text = ""
                    try:
                        first_page = pdf.pages[0]
                        sample_text = first_page.extract_text() or ""
                    except:
                        pass
                    
                    # OCR artifacts: lots of spaces, random characters, poor structure
                    if sample_text:
                        space_ratio = sample_text.count(' ') / len(sample_text) if len(sample_text) > 0 else 0
                        if space_ratio > 0.3:  # Too many spaces might indicate OCR
                            logger.debug("PDF appears scanned: high space ratio suggests OCR output")
                            return 'scanned'
                
                # Default: assume text-based if we can't determine (safer for text-based PDFs)
                logger.debug("Could not determine PDF type, defaulting to text-based")
                return 'text_based'
                
        except Exception as e:
            logger.warning(f"Error detecting PDF type: {str(e)}, defaulting to scanned")
            return 'scanned'  # Default to scanned if detection fails
    
    def _extract_from_pdf(self, file_path: Path) -> str:
        """
        Smart PDF extraction: routes to best method based on PDF type and quality.
        - Text-based PDFs: Uses pdfplumber (fast, local, free)
        - Scanned PDFs: Uses Google Document AI (primary) with pytesseract fallback
        
        Strategy:
        1. Try quick text extraction with pdfplumber
        2. If quality is poor (< 50 chars/page avg), treat as scanned and use Document AI
        3. Otherwise use pdfplumber result (fast and free)
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text content
        """
        # Smart routing: Check PDF structure to determine best extraction method
        try:
            with pdfplumber.open(file_path) as pdf:
                total_chars = 0
                pages_with_text = 0
                pages_with_text_objects = 0
                
                # Sample first 3 pages to check quality
                sample_pages = min(3, len(pdf.pages))
                has_cid_codes = False  # CID codes indicate scanned PDFs
                
                for page_num in range(sample_pages):
                    page = pdf.pages[page_num]
                    
                    # Check 1: Extract text
                    page_text = page.extract_text()
                    if page_text:
                        text_length = len(page_text.strip())
                        total_chars += text_length
                        if text_length > 50:
                            pages_with_text += 1
                        
                        # Check for CID codes (Character ID) - clear indicator of scanned PDF
                        if '(cid:' in page_text.lower() or page_text.count('(cid:') > 5:
                            has_cid_codes = True
                            logger.debug(f"Page {page_num + 1} contains CID codes - likely scanned PDF")
                    
                    # Check 2: Count text objects (most reliable indicator of text-based PDF)
                    try:
                        chars = page.chars
                        if chars and len(chars) > 100:  # Need substantial text objects
                            pages_with_text_objects += 1
                    except:
                        pass
                
                # Calculate metrics
                avg_text_per_page = total_chars / sample_pages if sample_pages > 0 else 0
                text_object_ratio = pages_with_text_objects / sample_pages if sample_pages > 0 else 0
                
                # Decision logic:
                # 1. If PDF has CID codes, it's definitely scanned (use Document AI)
                if has_cid_codes:
                    logger.info(f"PDF contains CID codes - detected as scanned PDF, using Document AI")
                    return self._extract_scanned_pdf(file_path)
                
                # 2. If most pages have text objects (>70%), it's definitely text-based
                if text_object_ratio > 0.7 and not has_cid_codes:
                    logger.info(f"PDF is text-based ({pages_with_text_objects}/{sample_pages} pages have text objects), using pdfplumber")
                    return self._extract_text_based_pdf(file_path)
                
                # 3. If no text objects and low text, it's scanned
                if text_object_ratio == 0 and avg_text_per_page < 100:
                    logger.info(f"PDF appears scanned (no text objects, avg {avg_text_per_page:.1f} chars/page), using Document AI")
                    return self._extract_scanned_pdf(file_path)
                
                # 3. If we got good average text (>200 chars/page), use pdfplumber
                if avg_text_per_page > 200:
                    logger.info(f"PDF appears text-based (avg {avg_text_per_page:.1f} chars/page), using pdfplumber")
                    return self._extract_text_based_pdf(file_path)
                
                # 4. Borderline: Try pdfplumber first, check quality, fallback if needed
                logger.info(f"PDF quality unclear (avg {avg_text_per_page:.1f} chars/page, {pages_with_text_objects}/{sample_pages} pages with text objects)")
                logger.info("Trying pdfplumber first, will fallback to Document AI if quality is poor")
                
                text = self._extract_text_based_pdf(file_path)
                
                # Quality check: if result is substantial and looks good, use it
                if text and len(text.strip()) > 1000:
                    # Check for OCR artifacts (too many spaces, poor structure)
                    space_ratio = text.count(' ') / len(text) if len(text) > 0 else 0
                    if space_ratio < 0.25:  # Reasonable space ratio
                        logger.info("pdfplumber extraction quality is good, using it")
                        return text
                
                # Fallback to Document AI for better quality
                logger.info("pdfplumber result quality insufficient, using Document AI")
                return self._extract_scanned_pdf(file_path)
                    
        except Exception as e:
            logger.warning(f"Error in PDF routing: {str(e)}, trying Document AI")
            return self._extract_scanned_pdf(file_path)
    
    def _extract_text_based_pdf(self, file_path: Path) -> str:
        """
        Extract text from text-based PDF using pdfplumber with table extraction.
        Fast, local, and free method for PDFs with selectable text.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text content
        """
        text_parts = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Strategy 1: Extract text directly (best for text-based PDFs)
                    page_text = page.extract_text(layout=True)
                    
                    # Strategy 2: If direct extraction is poor, try extracting tables and text separately
                    if not page_text or len(page_text.strip()) < 50:
                        # Try extracting tables as structured text
                        tables = page.extract_tables()
                        table_texts = []
                        for table in tables:
                            if table:
                                # Convert table to readable text format
                                table_rows = []
                                for row in table:
                                    if row:
                                        # Filter out None values and join cells
                                        clean_row = [str(cell).strip() if cell else "" for cell in row]
                                        row_text = " | ".join(clean_row)
                                        if row_text.strip():
                                            table_rows.append(row_text)
                                if table_rows:
                                    table_texts.append("\n".join(table_rows))
                        
                        if table_texts:
                            page_text = "\n\n".join(table_texts)
                    
                    # Strategy 3: Extract words/chars if text extraction fails
                    if not page_text or len(page_text.strip()) < 10:
                        words = page.extract_words()
                        if words:
                            # Group words by their y-coordinate to preserve line structure
                            lines = {}
                            for word in words:
                                y = round(word['top'])
                                if y not in lines:
                                    lines[y] = []
                                lines[y].append(word['text'])
                            
                            # Sort by y-coordinate and join words
                            sorted_lines = sorted(lines.items())
                            page_text = "\n".join([" ".join(words) for _, words in sorted_lines])
                    
                    if page_text:
                        # Clean PDF encoding artifacts
                        cleaned_text = self._clean_text(page_text)
                        if cleaned_text:
                            text_parts.append(cleaned_text)
                            logger.debug(f"Extracted {len(cleaned_text)} chars from page {page_num}")
            
            if text_parts:
                result = "\n\n".join(text_parts)
                logger.info(f"Extracted {len(result)} total characters from text-based PDF ({len(text_parts)} pages) using pdfplumber")
                return result
            else:
                logger.warning(f"No text extracted from text-based PDF {file_path}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from text-based PDF {file_path}: {str(e)}", exc_info=True)
            return ""
    
    def _extract_scanned_pdf(self, file_path: Path) -> str:
        """
        Extract text from scanned/image-based PDF.
        Primary: Google Document AI (high quality)
        Fallback: pytesseract OCR (local, free)
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text content
        """
        # Try Google Document AI first (if enabled and available)
        if settings.GOOGLE_DOCAI_ENABLED and DOCAI_AVAILABLE:
            try:
                logger.info(f"Attempting Google Document AI extraction for scanned PDF: {file_path.name}")
                text = self._extract_with_document_ai(file_path)
                if text and len(text.strip()) > 100:  # Only use if we got substantial text
                    logger.info(f"Successfully extracted {len(text)} characters using Google Document AI")
                    return text
                else:
                    logger.debug("Document AI extraction returned insufficient text, falling back to pytesseract OCR")
            except Exception as e:
                logger.warning(f"Document AI extraction failed: {str(e)}, falling back to pytesseract OCR")
        
        # Fallback to pytesseract OCR
        if OCR_AVAILABLE:
            logger.info(f"Using pytesseract OCR fallback for scanned PDF: {file_path.name}")
            return self._extract_with_ocr(file_path)
        else:
            logger.error("OCR not available and Document AI failed. Cannot extract text from scanned PDF.")
            return ""
    
    def _extract_from_word(self, file_path: Path, format_type: str = 'docx') -> str:
        """
        Extract text from Word document (.docx or .doc)
        
        Args:
            file_path: Path to Word document
            format_type: 'docx' or 'doc'
        """
        if format_type == 'docx':
            try:
                doc = DocxDocument(str(file_path))
                text_parts = []
                
                # Extract paragraphs
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)
                
                # Extract tables
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            table_text.append(row_text)
                    if table_text:
                        text_parts.append("\n".join(table_text))
                
                result = "\n".join(text_parts)
                logger.info(f"Extracted {len(result)} characters from Word document")
                return result
            except Exception as e:
                logger.error(f"Error extracting text from Word docx {file_path}: {str(e)}")
                return ""
        else:  # .doc format
            # For .doc files, we need to use textract or antiword
            # Since we want free libraries only, try to read as text or use OCR
            logger.warning(f".doc format not fully supported. Attempting basic extraction for {file_path.name}")
            try:
                # Try reading as binary and looking for text
                with open(file_path, 'rb') as f:
                    content = f.read()
                    # Try to extract readable text from binary
                    text = content.decode('utf-8', errors='ignore')
                    # Filter out non-printable characters
                    text = ''.join(char for char in text if char.isprintable() or char in '\n\r\t')
                    # Clean up using the standardized cleaner
                    text = self._clean_text(text, preserve_layout=True)
                    if len(text.strip()) > 100:  # Only return if we got substantial text
                        logger.info(f"Extracted {len(text)} characters from .doc file (basic extraction)")
                        return text
                    else:
                        # If basic extraction failed, try OCR
                        if OCR_AVAILABLE:
                            logger.info("Basic extraction failed, attempting OCR for .doc file")
                            return self._extract_with_ocr(file_path)
                        return ""
            except Exception as e:
                logger.error(f"Error extracting text from Word doc {file_path}: {str(e)}")
                if OCR_AVAILABLE:
                    return self._extract_with_ocr(file_path)
                return ""
    
    def _extract_from_excel(self, file_path: Path, format_type: str = 'xlsx') -> str:
        """
        Extract text from Excel file (.xlsx, .xls, or .csv)
        
        Args:
            file_path: Path to Excel file
            format_type: 'xlsx', 'xls', or 'csv'
        """
        if format_type == 'csv':
            if PANDAS_AVAILABLE and _pd is not None:
                try:
                    # Try different encodings
                    for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                        try:
                            df = _pd.read_csv(str(file_path), encoding=encoding, on_bad_lines='skip')
                            text = df.to_string(index=False)
                            logger.info(f"Extracted {len(text)} characters from CSV file")
                            return text
                        except UnicodeDecodeError:
                            continue
                    return ""
                except Exception as e:
                    logger.error(f"Error extracting text from CSV {file_path}: {str(e)}")
                    return ""
            else:
                # Fallback: read as text
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        return f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file_path, 'r', encoding='latin-1') as f:
                            return f.read()
                    except Exception as e:
                        logger.error(f"Error reading CSV {file_path}: {str(e)}")
                        return ""
        
        elif format_type == 'xls':
            if XLRD_AVAILABLE:
                try:
                    import xlrd
                    workbook = xlrd.open_workbook(str(file_path))
                    text_parts = []
                    for sheet_name in workbook.sheet_names():
                        sheet = workbook.sheet_by_name(sheet_name)
                        sheet_text = []
                        for row_idx in range(sheet.nrows):
                            row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                            row_text = " | ".join(row_values)
                            if row_text.strip():
                                sheet_text.append(row_text)
                        if sheet_text:
                            text_parts.append(f"\n--- Sheet: {sheet_name} ---\n" + "\n".join(sheet_text))
                    result = "\n\n".join(text_parts)
                    logger.info(f"Extracted {len(result)} characters from .xls file")
                    return result
                except Exception as e:
                    logger.error(f"Error extracting text from .xls {file_path}: {str(e)}")
                    return ""
            else:
                logger.warning("xlrd not available. Cannot extract from .xls files.")
                return ""
        
        else:  # .xlsx
            try:
                workbook = load_workbook(str(file_path), data_only=True)
                text_parts = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_text = []
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell else "" for cell in row])
                        if row_text.strip():
                            sheet_text.append(row_text)
                    if sheet_text:
                        text_parts.append(f"\n--- Sheet: {sheet_name} ---\n" + "\n".join(sheet_text))
                result = "\n\n".join(text_parts)
                logger.info(f"Extracted {len(result)} characters from Excel file")
                return result
            except Exception as e:
                logger.error(f"Error extracting text from Excel {file_path}: {str(e)}")
                return ""
    
    def _extract_from_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation (.pptx)"""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available. Cannot extract from PowerPoint files.")
            return ""
        
        try:
            if _Presentation is None:
                return ""
            prs = _Presentation(str(file_path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if hasattr(shape, "table"):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                            if row_text:
                                table_text.append(row_text)
                        if table_text:
                            slide_text.append("\n".join(table_text))
                
                if slide_text:
                    text_parts.append(f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            
            result = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(result)} characters from PowerPoint presentation")
            return result
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint {file_path}: {str(e)}")
            return ""
    
    def _extract_from_image(self, file_path: Path) -> str:
        """Extract text from image files using OCR"""
        if not OCR_AVAILABLE:
            logger.warning("OCR libraries not available. Cannot extract text from images.")
            return ""
        
        try:
            return self._extract_with_ocr(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {str(e)}")
            return ""
    
    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text file (.txt)"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    
                    # Skip the header if it's from our extraction (starts with "Source URL:")
                    if text.startswith("Source URL:"):
                        # Find where the actual content starts (after the separator line)
                        lines = text.split('\n')
                        content_start = 0
                        for i, line in enumerate(lines):
                            if line.startswith("=" * 80) or line.startswith("=" * 40):
                                content_start = i + 1
                                break
                        if content_start > 0:
                            text = '\n'.join(lines[content_start:])
                    
                    logger.info(f"Read {len(text)} characters from TXT file")
                    return text
                except UnicodeDecodeError:
                    continue
            
            logger.warning(f"Could not decode text file {file_path} with any encoding")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from text file {file_path}: {str(e)}")
            return ""
    
    def _extract_from_rtf(self, file_path: Path) -> str:
        """Extract text from RTF file"""
        if RTF_AVAILABLE and _rtf_to_text is not None:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    rtf_content = f.read()
                text = _rtf_to_text(rtf_content)
                logger.info(f"Extracted {len(text)} characters from RTF file")
                return text
            except Exception as e:
                logger.error(f"Error extracting text from RTF {file_path}: {str(e)}")
                return ""
        else:
            # Fallback: try to read as plain text (will have RTF codes but might work)
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                # Basic RTF cleanup - remove RTF control codes
                text = re.sub(r'\\[a-z]+\d*\s?', '', text)
                text = re.sub(r'\{[^}]*\}', '', text)
                text = re.sub(r'\s+', ' ', text)
                logger.info(f"Extracted {len(text)} characters from RTF file (basic extraction)")
                return text
            except Exception as e:
                logger.error(f"Error reading RTF file {file_path}: {str(e)}")
                return ""
    
    def _extract_from_markdown(self, file_path: Path) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            if MARKDOWN_AVAILABLE and _markdown is not None:
                # Convert markdown to plain text
                html = _markdown.markdown(md_content)
                # Remove HTML tags
                text = re.sub(r'<[^>]+>', '', html)
                logger.info(f"Extracted {len(text)} characters from Markdown file")
                return text
            else:
                # Fallback: return markdown as-is (still readable)
                logger.info(f"Read {len(md_content)} characters from Markdown file (raw)")
                return md_content
        except Exception as e:
            logger.error(f"Error extracting text from Markdown {file_path}: {str(e)}")
            return ""
    
    def _preprocess_image_for_ocr(self, image_path: Path) -> Optional[Any]:
        """
        Preprocess image file for better OCR results.
        Uses the same enhanced preprocessing as _preprocess_image_for_ocr_from_array.
        """
        if not OCR_AVAILABLE or _cv2 is None or _np is None:
            return None

        try:
            # Read image
            img = _cv2.imread(str(image_path))
            if img is None:
                logger.error(f"Could not read image {image_path}")
                return None

            # Convert BGR to RGB (OpenCV reads as BGR)
            img_rgb = _cv2.cvtColor(img, _cv2.COLOR_BGR2RGB)

            # Use the enhanced preprocessing
            return self._preprocess_image_for_ocr_from_array(img_rgb)
        except Exception as e:
            logger.warning(f"Image preprocessing failed: {str(e)}, using original image")
            try:
                img = _cv2.imread(str(image_path))
                if img is not None:
                    gray = _cv2.cvtColor(img, _cv2.COLOR_BGR2GRAY)
                    return self._preprocess_image_for_ocr_from_array(gray)
                return None
            except Exception:
                return None
    
    def _extract_with_ocr(self, file_path: Path) -> str:
        """
        Extract text using OCR (for scanned PDFs and images).
        Uses pytesseract with enhanced preprocessing and multiple OCR strategies.
        """
        if not OCR_AVAILABLE:
            logger.warning("OCR not available")
            return ""
        
        try:
            text_parts = []
            
            # Check if it's a PDF that needs OCR
            if file_path.suffix.lower() == '.pdf':
                # Convert PDF pages to images and OCR each
                try:
                    from pdf2image.pdf2image import convert_from_path
                    # Use higher DPI for better quality (300-400 is optimal)
                    images = convert_from_path(str(file_path), dpi=300, fmt='png')
                    logger.info(f"Converted PDF to {len(images)} images for OCR at 300 DPI")
                    
                    for page_num, image in enumerate(images, 1):
                        # Preprocess image
                        if _np is None or _cv2 is None:
                            continue
                        img_array = _np.array(image)
                        if len(img_array.shape) == 3:
                            gray = _cv2.cvtColor(img_array, _cv2.COLOR_RGB2GRAY)
                        else:
                            gray = img_array
                        
                        # Apply enhanced preprocessing
                        preprocessed = self._preprocess_image_for_ocr_from_array(gray)
                        
                        # Try multiple OCR strategies and use the best result
                        page_text = self._ocr_with_multiple_strategies(preprocessed, page_num)
                        
                        if page_text.strip():
                            # Clean OCR artifacts
                            cleaned_text = self._clean_ocr_text(page_text)
                            if cleaned_text:
                                text_parts.append(cleaned_text)
                                logger.info(f"OCR extracted {len(cleaned_text)} chars from page {page_num}")
                    
                    result = "\n\n".join(text_parts)
                    logger.info(f"OCR extracted {len(result)} total characters from PDF ({len(text_parts)} pages)")
                    return result
                except ImportError:
                    logger.warning("pdf2image not available. Cannot OCR PDF files.")
                    return ""
                except Exception as e:
                    logger.error(f"Error during PDF OCR: {str(e)}", exc_info=True)
                    return ""
            
            else:
                # It's an image file
                # Preprocess image
                preprocessed = self._preprocess_image_for_ocr(file_path)
                if preprocessed is None:
                    # Fallback: use original image
                    try:
                        if _PilImage is None or _np is None or _cv2 is None:
                            return ""
                        img = _PilImage.open(str(file_path))
                        img_array = _np.array(img)
                        if len(img_array.shape) == 3:
                            gray = _cv2.cvtColor(img_array, _cv2.COLOR_RGB2GRAY)
                        else:
                            gray = img_array
                        preprocessed = self._preprocess_image_for_ocr_from_array(gray)
                        text = self._ocr_with_multiple_strategies(preprocessed, 1)
                    except Exception as e:
                        logger.error(f"OCR failed: {str(e)}")
                        return ""
                else:
                    # Use preprocessed image with multiple strategies
                    text = self._ocr_with_multiple_strategies(preprocessed, 1)
                
                # Clean OCR artifacts
                cleaned_text = self._clean_ocr_text(text)
                logger.info(f"OCR extracted {len(cleaned_text)} characters from image")
                return cleaned_text
        
        except Exception as e:
            logger.error(f"Error during OCR extraction: {str(e)}", exc_info=True)
            return ""
    
    def _ocr_with_multiple_strategies(self, preprocessed_img: Any, page_num: int) -> str:
        """
        Try multiple OCR strategies (PSM modes) and return the best result.
        PSM (Page Segmentation Mode) modes:
        - 3: Fully automatic page segmentation, but no OSD (default)
        - 6: Assume a single uniform block of text
        - 11: Sparse text (find as much text as possible in no particular order)
        - 12: Sparse text with OSD
        """
        results = []
        
        if _pytesseract is None:
            return ""

        def _get_str(res: Any) -> str:
            if isinstance(res, dict):
                if 'text' in res:
                    text_val = res['text']
                    if isinstance(text_val, list):
                        return " ".join(str(t) for t in text_val if str(t).strip())
                    return str(text_val)
                return str(res)
            return str(res) if res is not None else ""

        # Strategy 1: Default PSM mode (3) - good for most documents
        try:
            config1 = '--psm 3 --oem 3 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz.,;:!?()[]{}\'"- /\\'
            text1 = _get_str(_pytesseract.image_to_string(preprocessed_img, lang='eng', config=config1))
            if text1.strip():
                results.append(('psm3', text1))
        except Exception as e:
            logger.debug(f"PSM 3 failed for page {page_num}: {str(e)}")

        # Strategy 2: Single block mode (6) - good for single column documents
        try:
            config2 = '--psm 6 --oem 3'
            text2 = _get_str(_pytesseract.image_to_string(preprocessed_img, lang='eng', config=config2))
            if text2.strip():
                results.append(('psm6', text2))
        except Exception as e:
            logger.debug(f"PSM 6 failed for page {page_num}: {str(e)}")

        # Strategy 3: Sparse text mode (11) - good for documents with scattered text
        try:
            config3 = '--psm 11 --oem 3'
            text3 = _get_str(_pytesseract.image_to_string(preprocessed_img, lang='eng', config=config3))
            if text3.strip():
                results.append(('psm11', text3))
        except Exception as e:
            logger.debug(f"PSM 11 failed for page {page_num}: {str(e)}")

        # Strategy 4: Auto PSM with OSD (12) - good for complex layouts
        try:
            config4 = '--psm 12 --oem 3'
            text4 = _get_str(_pytesseract.image_to_string(preprocessed_img, lang='eng', config=config4))
            if text4.strip():
                results.append(('psm12', text4))
        except Exception as e:
            logger.debug(f"PSM 12 failed for page {page_num}: {str(e)}")

        if not results:
            # Fallback: basic OCR without special config
            try:
                fallback_text = _get_str(_pytesseract.image_to_string(preprocessed_img, lang='eng'))
                return fallback_text
            except Exception:
                return ""
        
        # Choose the result with the most text (usually most accurate)
        # Also prefer results with more alphanumeric characters
        best_result = max(results, key=lambda x: (
            len(x[1]),
            sum(1 for c in x[1] if c.isalnum())
        ))
        
        logger.debug(f"Page {page_num}: Using {best_result[0]} OCR strategy ({len(best_result[1])} chars)")
        return best_result[1]
    
    def _clean_ocr_text(self, text: str) -> str:
        """
        Clean OCR artifacts and common errors.
        """
        if not text:
            return ""
        
        # Fix common OCR errors
        replacements = {
            # Common character misrecognitions
            r'[|]': 'I',  # | often misread as I
            r'[0O]': 'O',  # 0 vs O (context-dependent, but O is more common in text)
            r'rn': 'm',  # rn often misread as m
            r'vv': 'w',  # vv often misread as w
            r'ii': 'n',  # ii often misread as n
        }
        
        # Fix spacing issues
        text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single space
        
        # Fix line breaks (preserve intentional breaks)
        text = re.sub(r'\n\s*\n+', '\n\n', text)  # Multiple newlines to double
        
        # Remove common OCR artifacts
        text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)  # Remove non-printable except newlines/tabs
        
        # Fix common word breaks
        text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)  # Fix hyphenated words split across lines
        
        # Remove excessive punctuation
        text = re.sub(r'[.]{3,}', '...', text)  # Multiple periods to ellipsis
        
        return text.strip()
    
    def _preprocess_image_for_ocr_from_array(self, img_array: Any) -> Any:
        """
        Preprocess image array for OCR with enhanced techniques.
        Applies multiple preprocessing steps for better OCR accuracy.
        """
        if not OCR_AVAILABLE or _cv2 is None or _np is None:
            return img_array

        try:
            # Ensure image is grayscale
            if len(img_array.shape) == 3:
                gray = _cv2.cvtColor(img_array, _cv2.COLOR_RGB2GRAY)
            else:
                gray = img_array.copy()

            # Resize if image is too small (improves OCR accuracy)
            height, width = gray.shape
            if height < 300 or width < 300:
                scale = max(300 / height, 300 / width)
                new_width = int(width * scale)
                new_height = int(height * scale)
                gray = _cv2.resize(gray, (new_width, new_height), interpolation=_cv2.INTER_CUBIC)

            # Apply aggressive denoising (better for scanned documents)
            denoised = _cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)

            # Enhance contrast using CLAHE (adaptive histogram equalization)
            clahe = _cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8, 8))
            enhanced = clahe.apply(denoised)

            # Apply morphological operations to clean up the image
            kernel = _np.ones((2, 2), _np.uint8)
            enhanced = _cv2.morphologyEx(enhanced, _cv2.MORPH_CLOSE, kernel)

            # Deskew detection and correction
            enhanced = self._deskew_image(enhanced)

            # Apply adaptive thresholding for better text/background separation
            binary = _cv2.adaptiveThreshold(
                enhanced, 255, _cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                _cv2.THRESH_BINARY, 11, 2
            )

            # Alternative: Otsu's thresholding (sometimes better)
            _, binary_otsu = _cv2.threshold(enhanced, 0, 255, _cv2.THRESH_BINARY + _cv2.THRESH_OTSU)

            # Use the binary image with better contrast
            # Compare both and use the one with better contrast
            contrast_adaptive = _np.std(binary)
            contrast_otsu = _np.std(binary_otsu)

            final_binary = binary_otsu if contrast_otsu > contrast_adaptive else binary

            # Final cleanup: remove small noise
            kernel_clean = _np.ones((1, 1), _np.uint8)
            final_binary = _cv2.morphologyEx(final_binary, _cv2.MORPH_OPEN, kernel_clean)

            return final_binary
        except Exception as e:
            logger.warning(f"Image preprocessing failed: {str(e)}, using original image")
            return img_array
    
    @staticmethod
    def _deskew_image(image: Any) -> Any:
        """
        Detect and correct skew in scanned images.
        Returns deskewed image.
        """
        if _np is None or _cv2 is None:
            return image
        try:
            # Convert to binary for skew detection
            coords = _np.column_stack(_np.where(image > 0))
            if len(coords) == 0:
                return image

            # Find minimum area rectangle
            angle = _cv2.minAreaRect(coords)[-1]

            # Correct angle
            if angle < -45:
                angle = -(90 + angle)
            else:
                angle = -angle

            # Only correct if angle is significant (> 0.5 degrees)
            if abs(angle) > 0.5:
                (h, w) = image.shape[:2]
                center = (w // 2, h // 2)
                M = _cv2.getRotationMatrix2D(center, angle, 1.0)
                rotated = _cv2.warpAffine(image, M, (w, h), flags=_cv2.INTER_CUBIC,
                                        borderMode=_cv2.BORDER_REPLICATE)
                return rotated

            return image
        except Exception as e:
            logger.debug(f"Deskew failed: {str(e)}, returning original image")
            return image
